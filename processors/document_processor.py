import os
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation

class DocumentProcessor:
    """Handles processing of various document formats"""
    
    def process(self, file_path: str) -> str:
        """
        Process a document file and extract text content
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text content
        """
        try:
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension == '.pdf':
                return self._process_pdf(file_path)
            elif file_extension == '.docx':
                return self._process_docx(file_path)
            elif file_extension == '.pptx':
                return self._process_pptx(file_path)
            elif file_extension in ['.md', '.txt']:
                return self._process_text(file_path)
            else:
                raise ValueError(f"Unsupported document format: {file_extension}")
                
        except Exception as e:
            raise Exception(f"Error processing document {file_path}: {str(e)}")
    
    def _process_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
        
        return text.strip()
    
    def _process_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                text += "\n"
                        
        except Exception as e:
            raise Exception(f"Error reading DOCX: {str(e)}")
        
        return text.strip()
    
    def _process_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                text += "\n"
                        
        except Exception as e:
            raise Exception(f"Error reading PPTX: {str(e)}")
        
        return text.strip()
    
    def _process_text(self, file_path: str) -> str:
        """Extract text from plain text or markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read().strip()
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read().strip()
        except Exception as e:
            raise Exception(f"Error reading text file: {str(e)}")
